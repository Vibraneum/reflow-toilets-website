#!/usr/bin/env python3
"""
Extract text and images from PDFs and PPTs for ReFlow Toilets website.
"""

import os
import sys
import json
from pathlib import Path
from typing import Dict, List, Any

try:
    import PyPDF2
    from pdf2image import convert_from_path
    from PIL import Image
except ImportError:
    print("Installing required packages...")
    os.system(f"{sys.executable} -m pip install PyPDF2 pdf2image pillow")
    import PyPDF2
    from pdf2image import convert_from_path
    from PIL import Image

# Add parent directory to path
sys.path.insert(0, str(Path(__file__).parent.parent))

def extract_text_from_pdf(pdf_path: str) -> Dict[str, Any]:
    """Extract text content from PDF."""
    content = {
        "file": os.path.basename(pdf_path),
        "pages": [],
        "full_text": ""
    }
    
    try:
        with open(pdf_path, 'rb') as file:
            pdf_reader = PyPDF2.PdfReader(file)
            num_pages = len(pdf_reader.pages)
            
            for page_num in range(num_pages):
                page = pdf_reader.pages[page_num]
                text = page.extract_text()
                content["pages"].append({
                    "page_number": page_num + 1,
                    "text": text
                })
                content["full_text"] += f"\n\n--- Page {page_num + 1} ---\n\n{text}"
    
    except Exception as e:
        print(f"Error extracting text from {pdf_path}: {e}")
        content["error"] = str(e)
    
    return content

def extract_images_from_pdf(pdf_path: str, output_dir: str) -> List[str]:
    """Extract images from PDF."""
    image_paths = []
    
    try:
        # Convert PDF pages to images
        images = convert_from_path(pdf_path, dpi=200)
        
        base_name = Path(pdf_path).stem
        os.makedirs(output_dir, exist_ok=True)
        
        for i, image in enumerate(images):
            image_path = os.path.join(output_dir, f"{base_name}_page_{i+1}.png")
            image.save(image_path, "PNG")
            image_paths.append(image_path)
            print(f"Extracted image: {image_path}")
    
    except Exception as e:
        print(f"Error extracting images from {pdf_path}: {e}")
        # Fallback: try using pdf2image with poppler
        try:
            images = convert_from_path(pdf_path, dpi=150)
            base_name = Path(pdf_path).stem
            os.makedirs(output_dir, exist_ok=True)
            
            for i, image in enumerate(images):
                image_path = os.path.join(output_dir, f"{base_name}_page_{i+1}.png")
                image.save(image_path, "PNG")
                image_paths.append(image_path)
        except Exception as e2:
            print(f"Fallback also failed: {e2}")
    
    return image_paths

def process_powerpoint(pptx_path: str, output_dir: str) -> Dict[str, Any]:
    """Extract content from PowerPoint file."""
    content = {
        "file": os.path.basename(pptx_path),
        "slides": [],
        "full_text": ""
    }
    
    try:
        from pptx import Presentation
        
        prs = Presentation(pptx_path)
        
        for slide_num, slide in enumerate(prs.slides):
            slide_text = []
            
            for shape in slide.shapes:
                if hasattr(shape, "text"):
                    slide_text.append(shape.text)
            
            slide_content = "\n".join(slide_text)
            content["slides"].append({
                "slide_number": slide_num + 1,
                "text": slide_content
            })
            content["full_text"] += f"\n\n--- Slide {slide_num + 1} ---\n\n{slide_content}"
    
    except ImportError:
        print("python-pptx not installed. Installing...")
        os.system(f"{sys.executable} -m pip install python-pptx")
        return process_powerpoint(pptx_path, output_dir)
    except Exception as e:
        print(f"Error processing PowerPoint {pptx_path}: {e}")
        content["error"] = str(e)
    
    return content

def main():
    """Main extraction function."""
    # Define source PDFs and PPTs
    base_dir = Path(__file__).parent.parent.parent
    source_files = [
        base_dir / "bcrt1.pdf",
        base_dir / "Loocafe - Reflow.pdf",
        base_dir / "ReFlow Toilets" / "ReFlow Toilets" / "info for ai reflow_compressed.pdf (Presentation).pdf",
        base_dir / "ReFlow Toilets" / "ReFlow Toilets Website assets" / "iMAGES irise toilet.pdf",
        base_dir / "ReFlow Toilets" / "ReFlow Toilets Website assets" / "tech reflow.pdf",
        base_dir / "Eco-san company and product introduction_CRT&HRT version (1).pptx",
    ]
    
    # Output directories
    output_base = Path(__file__).parent.parent / "extracted-content"
    text_output = output_base / "text"
    images_output = output_base / "images"
    json_output = output_base / "json"
    
    os.makedirs(text_output, exist_ok=True)
    os.makedirs(images_output, exist_ok=True)
    os.makedirs(json_output, exist_ok=True)
    
    all_content = {}
    
    for source_file in source_files:
        if not source_file.exists():
            print(f"File not found: {source_file}")
            continue
        
        print(f"\nProcessing: {source_file.name}")
        
        file_ext = source_file.suffix.lower()
        base_name = source_file.stem
        
        if file_ext == ".pdf":
            # Extract text
            text_content = extract_text_from_pdf(str(source_file))
            
            # Save text
            text_file = text_output / f"{base_name}.txt"
            with open(text_file, "w", encoding="utf-8") as f:
                f.write(text_content["full_text"])
            
            # Extract images
            image_dir = images_output / base_name
            image_paths = extract_images_from_pdf(str(source_file), str(image_dir))
            text_content["images"] = image_paths
            
            # Save JSON
            json_file = json_output / f"{base_name}.json"
            with open(json_file, "w", encoding="utf-8") as f:
                json.dump(text_content, f, indent=2, ensure_ascii=False)
            
            all_content[base_name] = text_content
            
        elif file_ext in [".ppt", ".pptx"]:
            # Process PowerPoint
            ppt_content = process_powerpoint(str(source_file), str(images_output / base_name))
            
            # Save text
            text_file = text_output / f"{base_name}.txt"
            with open(text_file, "w", encoding="utf-8") as f:
                f.write(ppt_content["full_text"])
            
            # Save JSON
            json_file = json_output / f"{base_name}.json"
            with open(json_file, "w", encoding="utf-8") as f:
                json.dump(ppt_content, f, indent=2, ensure_ascii=False)
            
            all_content[base_name] = ppt_content
    
    # Save master JSON
    master_json = json_output / "all_content.json"
    with open(master_json, "w", encoding="utf-8") as f:
        json.dump(all_content, f, indent=2, ensure_ascii=False)
    
    print(f"\nExtraction complete!")
    print(f"Text files: {text_output}")
    print(f"Images: {images_output}")
    print(f"JSON data: {json_output}")
    print(f"Master JSON: {master_json}")

if __name__ == "__main__":
    main()

